#!/usr/bin/env python3
"""Generate sample PPT file for i-Write demo - Keynote compatible"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from lxml import etree
import os

def set_run_font(run, size_pt, bold=False, color_rgb=None):
    """Explicitly set font properties on the run element"""
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    if rPr is None:
        rPr = etree.SubElement(run._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    rPr.set('sz', str(int(size_pt * 100)))
    if bold:
        rPr.set('b', '1')
    # Remove any existing fill
    for fill in rPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill'):
        rPr.remove(fill)
    # Add explicit color
    if color_rgb:
        solidFill = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgbClr.set('val', color_rgb)
    # Add font references
    for tag in ['latin', 'ea', 'cs']:
        existing = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}' + tag)
        if existing is None:
            elem = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}' + tag)
            if tag == 'latin':
                elem.set('typeface', 'Helvetica Neue')
            else:
                elem.set('typeface', 'PingFang SC')

def add_text_run(tf, text, size_pt, bold=False, color='000000', align=PP_ALIGN.LEFT):
    """Add a paragraph with a run that has explicit font properties"""
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size_pt, bold, color)
    return p

def create_ppt(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide1.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True

    # First paragraph
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = "i-Write 项目周报"
    set_run_font(run1, 44, True, '000000')

    add_text_run(tf, "2026-06-16 ~ 2026-06-20", 20, False, '555555', PP_ALIGN.CENTER)
    add_text_run(tf, "陈强 · 技术负责人 · Nexora Tech", 16, False, '888888', PP_ALIGN.CENTER)

    # ── Slide 2: Goals ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "本周目标"
    set_run_font(run2, 32, True, '000000')

    for goal in ["1. 完成用户认证模块（OAuth2 + JWT）", "2. 修复 3 个高优 Bug（BUG-201/205/210）", "3. 完成支付系统技术方案设计"]:
        add_text_run(tf2, goal, 22, False, '333333')

    # ── Slide 3: Status Table ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "完成情况"
    set_run_font(run3, 32, True, '000000')

    # Table
    table = slide3.shapes.add_table(9, 3, Inches(1), Inches(1.5), Inches(11), Inches(5)).table

    # Header
    for i, h in enumerate(["任务", "状态", "负责人"]):
        cell = table.cell(0, i)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        set_run_font(run, 14, True, '000000')

    # Data
    data = [
        ["认证模块 OAuth2", "✅ 完成", "陈强"],
        ["Microsoft OAuth2", "✅ 完成", "陈强、刘伟"],
        ["GitHub OAuth", "✅ 完成", "刘伟"],
        ["登录页面 UI", "✅ 完成", "赵丽"],
        ["BUG-201/205/210", "✅ 完成", "赵丽、刘伟"],
        ["E2E 测试", "🔄 80%", "杨飞"],
        ["支付系统方案", "✅ 完成", "陈强、刘伟"],
        ["Token 刷新修复", "✅ 完成", "刘伟"],
    ]
    for r, row in enumerate(data, 1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            set_run_font(run, 14, False, '333333')

    # ── Slide 4: Sprint 进度（柱状图）──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "Sprint 进度"
    set_run_font(run4, 32, True, '000000')

    # Sprint progress bar chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Sprint 1', 'Sprint 2', 'Sprint 3', 'Sprint 4']
    chart_data.add_series('计划任务', (42, 51, 47, 55))
    chart_data.add_series('完成任务', (42, 51, 42, 55))
    chart = slide4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(11), Inches(5),
        chart_data
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # ── Slide 5: Issues ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    run5 = p5.add_run()
    run5.text = "遇到的问题与解决方案"
    set_run_font(run5, 32, True, '000000')

    txBox5b = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
    tf5b = txBox5b.text_frame
    tf5b.word_wrap = True

    issues = [
        ("问题 1: Azure AD redirect_uri 配置", True, 'CC0000'),
        ("→ 一个 App Registration 配多个 uri，用环境变量区分", False, '008800'),
        ("", False, '000000'),
        ("问题 2: Token 刷新竞态条件", True, 'CC0000'),
        ("→ 队列机制：刷新中请求排队等待，刷新后用新 Token 重发", False, '008800'),
        ("", False, '000000'),
        ("问题 3: GitHub OAuth 回调超时", True, 'CC0000'),
        ("→ 增加超时到 10s + 重试机制", False, '008800'),
    ]
    for text, bold, color in issues:
        add_text_run(tf5b, text, 18, bold, color)

    # ── Slide 6: Next Week ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    run6 = p6.add_run()
    run6.text = "下周计划"
    set_run_font(run6, 32, True, '000000')

    txBox6b = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf6b = txBox6b.text_frame
    tf6b.word_wrap = True

    for plan in ["1. 完成 E2E 测试并修复发现的问题", "2. 支付系统设计评审", "3. 开始支付模块开发", "4. 更新技术文档"]:
        add_text_run(tf6b, plan, 22, False, '333333')

    # ── Slide 7: RAG 引擎进展 ──
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox7 = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf7 = txBox7.text_frame
    p7 = tf7.paragraphs[0]
    run7 = p7.add_run()
    run7.text = "RAG 引擎进展"
    set_run_font(run7, 32, True, '000000')

    txBox7b = slide7.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf7b = txBox7b.text_frame
    tf7b.word_wrap = True
    rag_items = [
        "Query Expansion → Hybrid Search → Reranker → Groundedness Check",
        "",
        "✅ BM25 检索（MiniSearch + jieba-wasm）— Recall@10: 0.78",
        "✅ 向量检索（SiliconFlow bge-m3, 1024 维）— Recall@10: 0.82",
        "✅ RRF 融合（k=60）— Recall@10: 0.91",
        "✅ Reranker 三级降级（远程 API → 本地 Cross-Encoder → 启发式）",
        "✅ Groundedness Check — 准确率 92%",
        "",
        "端到端延迟: < 3s（P95）",
    ]
    for item in rag_items:
        add_text_run(tf7b, item, 18, False, '333333')

    # ── Slide 8: 认证模块架构 ──
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox8 = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf8 = txBox8.text_frame
    p8 = tf8.paragraphs[0]
    run8 = p8.add_run()
    run8.text = "认证模块架构"
    set_run_font(run8, 32, True, '000000')

    txBox8b = slide8.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf8b = txBox8b.text_frame
    tf8b.word_wrap = True
    auth_items = [
        "User → OAuth2 Provider → Callback → JWT Generation → Frontend Storage",
        "",
        "✅ Microsoft OAuth2（MSAL.js + PKCE）",
        "✅ GitHub OAuth（passport.js）",
        "✅ JWT Token 管理（Access Token 1h + Refresh Token 7d）",
        "✅ Token 刷新队列机制（修复 BUG-215 竞态条件）",
        "",
        "登录成功率: 99.8% | Token 刷新成功率: 99.95%",
    ]
    for item in auth_items:
        add_text_run(tf8b, item, 18, False, '333333')

    # ── Slide 9: Bug 趋势（柱状图）──
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox9 = slide9.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf9 = txBox9.text_frame
    p9 = tf9.paragraphs[0]
    run9 = p9.add_run()
    run9.text = "Bug 趋势"
    set_run_font(run9, 32, True, '000000')

    # Bug trend bar chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Sprint 1', 'Sprint 2', 'Sprint 3', 'Sprint 4']
    chart_data.add_series('新增 Bug', (3, 5, 4, 3))
    chart_data.add_series('修复 Bug', (0, 3, 5, 4))
    chart = slide9.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(11), Inches(5),
        chart_data
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # ── Slide 10: 遇到的问题与解决方案 ──
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox10 = slide10.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf10 = txBox10.text_frame
    p10 = tf10.paragraphs[0]
    run10 = p10.add_run()
    run10.text = "遇到的问题与解决方案"
    set_run_font(run10, 32, True, '000000')

    txBox10b = slide10.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
    tf10b = txBox10b.text_frame
    tf10b.word_wrap = True
    issues = [
        ("问题 1: Azure AD redirect_uri 配置", True, 'CC0000'),
        ("→ 一个 App Registration 配多个 uri，用环境变量区分", False, '008800'),
        ("", False, '000000'),
        ("问题 2: Token 刷新竞态条件（BUG-215）", True, 'CC0000'),
        ("→ 队列机制：刷新中请求排队等待，刷新后用新 Token 重发", False, '008800'),
        ("", False, '000000'),
        ("问题 3: GitHub OAuth 回调超时（BUG-216）", True, 'CC0000'),
        ("→ 增加超时到 10s + 重试机制", False, '008800'),
        ("", False, '000000'),
        ("问题 4: Gemini systemPromptMode 不兼容", True, 'CC0000'),
        ("→ ModelCapabilities 查询机制，自动适配参数差异", False, '008800'),
    ]
    for text, bold, color in issues:
        add_text_run(tf10b, text, 18, bold, color)

    # ── Slide 11: 团队协作 ──
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox11 = slide11.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf11 = txBox11.text_frame
    p11 = tf11.paragraphs[0]
    run11 = p11.add_run()
    run11.text = "团队协作"
    set_run_font(run11, 32, True, '000000')

    txBox11b = slide11.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf11b = txBox11b.text_frame
    tf11b.word_wrap = True
    team_items = [
        "技术部（8 人）：陈强、刘伟、赵丽、孙娜、王超、周敏、徐骏、杨飞",
        "产品部（2 人）：苏楠、黄薇",
        "设计部（2 人）：罗茜、何成",
        "市场/销售（2 人）：王莉、张伟",
        "法务（1 人）：唐敏",
        "客户成功（1 人）：李鑫",
        "管理层（3 人）：陈宇、王琳、赵军",
        "",
        "月度 PR: 195 | 测试覆盖率: 68% → 87%",
    ]
    for item in team_items:
        add_text_run(tf11b, item, 18, False, '333333')

    # ── Slide 12: Q&A ──
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox12 = slide12.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf12 = txBox12.text_frame
    p12 = tf12.paragraphs[0]
    p12.alignment = PP_ALIGN.CENTER
    run12 = p12.add_run()
    run12.text = "Q & A"
    set_run_font(run12, 44, True, '000000')
    add_text_run(tf12, "感谢关注 i-Write 项目进展", 20, False, '555555', PP_ALIGN.CENTER)

    # Save
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"  ✅ {os.path.basename(output_path)}")


def create_roadmap_ppt(output_path):
    """产品路线图 PPT（10 slides）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide1.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = "i-Write 产品路线图"
    set_run_font(run1, 44, True, '000000')
    add_text_run(tf, "Q3 2026 · Nexora Tech", 20, False, '555555', PP_ALIGN.CENTER)

    # ── Slide 2: 产品愿景 ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "产品愿景"
    set_run_font(run2, 32, True, '000000')
    txBox2b = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf2b = txBox2b.text_frame
    tf2b.word_wrap = True
    for item in [
        "成为企业级可信文档生成平台",
        "",
        "核心价值主张：让每一份文档都有据可查、有源可溯",
        "",
        "差异化优势：",
        "1. 溯源能力：每段文字都能追溯到知识库来源",
        "2. 知识库集成：支持 10+ 知识源类型",
        "3. Trust Score：5 维度评估文档可信度",
    ]:
        add_text_run(tf2b, item, 20, False, '333333')

    # ── Slide 3: 市场机会（饼图）──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "市场机会"
    set_run_font(run3, 32, True, '000000')

    # Market size pie chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Jasper', 'Copy.ai', 'Notion AI', '其他', 'i-Write 目标']
    chart_data.add_series('市场份额', (35, 25, 20, 10, 10))
    chart = slide3.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(1), Inches(1.5), Inches(5), Inches(5),
        chart_data
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # Market data text
    txBox3b = slide3.shapes.add_textbox(Inches(7), Inches(2), Inches(5), Inches(4))
    tf3b = txBox3b.text_frame
    tf3b.word_wrap = True
    for item in [
        "全球市场规模：$4.2B（2026）",
        "年增长率：28%",
        "i-Write 目标：$420M",
        "",
        "竞品定价：",
        "Jasper: $49/月",
        "Copy.ai: $36/月",
        "i-Write: $19/月",
    ]:
        add_text_run(tf3b, item, 18, False, '333333')

    # ── Slide 4: 核心功能 ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "核心功能"
    set_run_font(run4, 32, True, '000000')
    txBox4b = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf4b = txBox4b.text_frame
    tf4b.word_wrap = True
    for item in [
        "1. 知识库管理 — 10+ 知识源类型，自动切片 + Embedding",
        "2. RAG 引擎 — BM25 + 向量检索 + RRF 融合 + Reranker",
        "3. 可信生成 — 叙事引擎 + Groundedness Check",
        "4. 生成树溯源 — 每段文字追溯到 chunk 级别",
        "5. 评估体系 — Trust Metrics 5 维度评分",
        "6. 文档导出 — Word / PPT / Excel 三种格式",
    ]:
        add_text_run(tf4b, item, 20, False, '333333')

    # ── Slide 5: 技术路线 ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    run5 = p5.add_run()
    run5.text = "技术路线（5 Phase）"
    set_run_font(run5, 32, True, '000000')
    txBox5b = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf5b = txBox5b.text_frame
    tf5b.word_wrap = True
    for item in [
        "Phase 1 (6月): 基础设施 — 项目骨架、Provider、CI/CD ✅",
        "Phase 2 (6-7月): 知识管理 — 文件上传、知识库、Embedding ✅",
        "Phase 3 (7月): RAG 引擎 — BM25、向量检索、Reranker ✅",
        "Phase 4 (8月): 文档生成 — Word/PPT/Excel、叙事引擎 ✅",
        "Phase 5 (9月): 评估体系 — 在线评估、离线评估、历史对比 🔄",
    ]:
        add_text_run(tf5b, item, 20, False, '333333')

    # ── Slide 6: 竞品对比 ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    run6 = p6.add_run()
    run6.text = "竞品对比"
    set_run_font(run6, 32, True, '000000')
    table6 = slide6.shapes.add_table(5, 5, Inches(1), Inches(1.5), Inches(11), Inches(4)).table
    headers = ["功能", "i-Write", "Jasper", "Copy.ai", "Notion AI"]
    for i, h in enumerate(headers):
        cell = table6.cell(0, i)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        set_run_font(run, 14, True, '000000')
    data = [
        ["溯源能力", "✅ chunk 级别", "❌", "❌", "❌"],
        ["知识库集成", "✅ 10+ 源", "❌", "❌", "⚠️ 基础"],
        ["Trust Score", "✅ 5 维度", "❌", "❌", "❌"],
        ["价格", "$19/月", "$49/月", "$36/月", "$10/月"],
    ]
    for r, row in enumerate(data, 1):
        for c, text in enumerate(row):
            cell = table6.cell(r, c)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            set_run_font(run, 14, False, '333333')

    # ── Slide 7: 定价策略 ──
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox7 = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf7 = txBox7.text_frame
    p7 = tf7.paragraphs[0]
    run7 = p7.add_run()
    run7.text = "定价策略"
    set_run_font(run7, 32, True, '000000')
    plans = [("Free", "$0/月", "5 次生成/月", 1.5), ("Pro", "$19/月", "无限生成", 5), ("Team", "$49/月/人", "团队协作", 8.5)]
    for name, price, desc, left in plans:
        txBox = slide7.shapes.add_textbox(Inches(left), Inches(2), Inches(2.5), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        set_run_font(run, 28, True, '1a56db')
        txBox2 = slide7.shapes.add_textbox(Inches(left), Inches(2.8), Inches(2.5), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = price
        set_run_font(run2, 24, True, '333333')
        txBox3 = slide7.shapes.add_textbox(Inches(left), Inches(3.4), Inches(2.5), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = desc
        set_run_font(run3, 16, False, '666666')

    # ── Slide 8: GoToMarket 计划 ──
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox8 = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf8 = txBox8.text_frame
    p8 = tf8.paragraphs[0]
    run8 = p8.add_run()
    run8.text = "GoToMarket 计划"
    set_run_font(run8, 32, True, '000000')
    txBox8b = slide8.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf8b = txBox8b.text_frame
    tf8b.word_wrap = True
    for item in [
        "PLG（Product-Led Growth）：免费版 → 自然升级",
        "企业直销：目标 Q3 签约 10 家企业客户",
        "内容营销：技术博客 + 社交媒体（掘金、知乎）",
        "",
        "关键里程碑：",
        "7/15: Beta 发布 | 8/1: Landing page | 8/15: 首个付费客户",
    ]:
        add_text_run(tf8b, item, 20, False, '333333')

    # ── Slide 9: 团队与资源 ──
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox9 = slide9.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf9 = txBox9.text_frame
    p9 = tf9.paragraphs[0]
    run9 = p9.add_run()
    run9.text = "团队与资源"
    set_run_font(run9, 32, True, '000000')
    txBox9b = slide9.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf9b = txBox9b.text_frame
    tf9b.word_wrap = True
    for item in [
        "当前团队：18 人（技术 8 + 产品 2 + 设计 2 + 市场 2 + 法务 1 + CS 1 + 管理 3）",
        "Q3 末目标：25 人（+3 工程师、+2 销售、+2 CS）",
        "",
        "融资需求：Series A $5M",
        "用途：团队扩展（$2M）+ 市场推广（$1.5M）+ 基础设施（$1.5M）",
    ]:
        add_text_run(tf9b, item, 20, False, '333333')

    # ── Slide 10: 下一步行动 ──
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox10 = slide10.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf10 = txBox10.text_frame
    p10 = tf10.paragraphs[0]
    run10 = p10.add_run()
    run10.text = "下一步行动"
    set_run_font(run10, 32, True, '000000')
    txBox10b = slide10.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf10b = txBox10b.text_frame
    tf10b.word_wrap = True
    for item in [
        "1. 7 月：Beta 版本发布 + Landing page 上线",
        "2. 8 月：PLG 启动 + 企业直销启动",
        "3. 9 月：首批企业客户上线 + 产品发布会",
        "",
        "关键目标：Q3 末 DAU 1000 + 10 家企业客户",
    ]:
        add_text_run(tf10b, item, 20, False, '333333')

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"  ✅ {os.path.basename(output_path)}")


def create_investor_ppt(output_path):
    """投资人更新 PPT（8 slides）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide1.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = "Nexora Tech · 投资人更新"
    set_run_font(run1, 44, True, '000000')
    add_text_run(tf, "2026 年 6 月", 20, False, '555555', PP_ALIGN.CENTER)

    # ── Slide 2: 公司概览 ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "公司概览"
    set_run_font(run2, 32, True, '000000')
    txBox2b = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf2b = txBox2b.text_frame
    tf2b.word_wrap = True
    for item in [
        "公司：Nexora Tech",
        "产品：i-Write — 企业级可信文档生成平台",
        "团队：18 人（技术 8 + 产品 2 + 设计 2 + 市场 2 + 法务 1 + CS 1 + 管理 3）",
        "阶段：Alpha 版本就绪，准备 Beta 发布",
        "",
        "核心能力：溯源 + 知识库 + Trust Score",
    ]:
        add_text_run(tf2b, item, 20, False, '333333')

    # ── Slide 3: 产品进展 ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "产品进展（6 月）"
    set_run_font(run3, 32, True, '000000')
    txBox3b = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf3b = txBox3b.text_frame
    tf3b.word_wrap = True
    for item in [
        "✅ Sprint 1: 基础设施搭建（42 PR merged）",
        "✅ Sprint 2: 知识管理模块（51 PR merged）",
        "✅ Sprint 3: RAG 引擎 + 认证模块（47 PR merged）",
        "✅ Sprint 4: 文档生成 + 评估体系（55 PR merged）",
        "",
        "月度总计：195 PR merged | 测试覆盖率：68% → 87%",
        "i-Write Alpha 版本就绪",
    ]:
        add_text_run(tf3b, item, 20, False, '333333')

    # ── Slide 4: 市场数据（用户增长柱状图）──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "市场数据"
    set_run_font(run4, 32, True, '000000')

    # User growth bar chart
    chart_data = CategoryChartData()
    chart_data.categories = ['7 月', '8 月', '9 月', '10 月', '11 月', '12 月']
    chart_data.add_series('注册用户', (100, 500, 1500, 3000, 5000, 8000))
    chart_data.add_series('付费用户', (5, 25, 75, 150, 250, 400))
    chart = slide4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(6), Inches(5),
        chart_data
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # Market data text
    txBox4b = slide4.shapes.add_textbox(Inches(8), Inches(2), Inches(4), Inches(4))
    tf4b = txBox4b.text_frame
    tf4b.word_wrap = True
    for item in [
        "NPS: 8.0/10",
        "意向客户：",
        "Acme Corp",
        "200 人",
        "$117,600/年",
    ]:
        add_text_run(tf4b, item, 18, False, '333333')

    # ── Slide 5: 技术优势 ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    run5 = p5.add_run()
    run5.text = "技术优势"
    set_run_font(run5, 32, True, '000000')
    txBox5b = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf5b = txBox5b.text_frame
    tf5b.word_wrap = True
    for item in [
        "1. 溯源能力：chunk 级别溯源，每段文字有据可查",
        "2. RAG 引擎：Hybrid Search + RRF 融合，Recall@10: 0.91",
        "3. Trust Score：5 维度评估，准确率 92%",
        "",
        "技术壁垒：",
        "- 多模型自适应（ModelCapabilities 查询机制）",
        "- Reranker 三级降级（远程 → 本地 → 启发式）",
        "- Groundedness Check（句子级验证）",
    ]:
        add_text_run(tf5b, item, 20, False, '333333')

    # ── Slide 6: 商业模式 ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    run6 = p6.add_run()
    run6.text = "商业模式"
    set_run_font(run6, 32, True, '000000')
    txBox6b = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf6b = txBox6b.text_frame
    tf6b.word_wrap = True
    for item in [
        "SaaS 订阅模式：Free / Pro ($19) / Team ($49)",
        "年付优惠：8 折",
        "",
        "收入预测（Q3 末）：",
        "- 10 家企业客户 × $117,600/年 = $1.176M ARR",
        "- 500 个人用户 × $19/月 × 12 = $114K ARR",
        "- 总计：$1.29M ARR",
    ]:
        add_text_run(tf6b, item, 20, False, '333333')

    # ── Slide 7: 财务预测（柱状图）──
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox7 = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf7 = txBox7.text_frame
    p7 = tf7.paragraphs[0]
    run7 = p7.add_run()
    run7.text = "财务预测"
    set_run_font(run7, 32, True, '000000')

    # Revenue forecast bar chart
    chart_data = CategoryChartData()
    chart_data.categories = ['2026 Q3', '2026 Q4', '2027 Q1', '2027 Q2']
    chart_data.add_series('ARR ($M)', (1.29, 3.5, 6.5, 10))
    chart = slide7.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(6), Inches(5),
        chart_data
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # Financial text
    txBox7b = slide7.shapes.add_textbox(Inches(8), Inches(2), Inches(4), Inches(4))
    tf7b = txBox7b.text_frame
    tf7b.word_wrap = True
    for item in [
        "烧钱率：$150K/月",
        "Runway：18 个月",
        "",
        "盈亏平衡：",
        "2027 Q4（预计）",
    ]:
        add_text_run(tf7b, item, 18, False, '333333')

    # ── Slide 8: 融资需求 ──
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox8 = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf8 = txBox8.text_frame
    p8 = tf8.paragraphs[0]
    run8 = p8.add_run()
    run8.text = "融资需求"
    set_run_font(run8, 32, True, '000000')
    txBox8b = slide8.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf8b = txBox8b.text_frame
    tf8b.word_wrap = True
    for item in [
        "Series A: $5M",
        "",
        "用途：",
        "- 团队扩展：$2M（+7 人，18 → 25 人）",
        "- 市场推广：$1.5M（PLG + 企业直销）",
        "- 基础设施：$1.5M（服务器、LLM API、安全合规）",
        "",
        "目标：2027 Q2 达到 $10M ARR",
    ]:
        add_text_run(tf8b, item, 20, False, '333333')

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"  ✅ {os.path.basename(output_path)}")

if __name__ == "__main__":
    output_dir = os.path.join(os.getcwd(), "samples", "presentations")

    print("📽️ 生成 PPT 文件...\n")

    # 1. 项目周报（12 slides）
    create_ppt(os.path.join(output_dir, "项目周报-2026-06-20.pptx"))

    # 2. 产品路线图（10 slides）
    create_roadmap_ppt(os.path.join(output_dir, "产品路线图-Q3-2026.pptx"))

    # 3. 投资人更新（8 slides）
    create_investor_ppt(os.path.join(output_dir, "投资人更新-2026-06.pptx"))

    print("\n🎉 PPT 全部生成完成！")
