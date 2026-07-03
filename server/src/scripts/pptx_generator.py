#!/usr/bin/env python3
"""
PPTX Generator — 从 JSON 数据生成 PowerPoint 演示文稿。

用法: python3 pptx_generator.py <input.json> <output.pptx>

JSON 格式:
{
  "title": "演示文稿标题",
  "sections": [
    {
      "title": "章节标题",
      "content": "HTML 正文",
      "tables": [[["A","B"],["1","2"]], ...],
      "chartSpecs": [{"type":"column","title":"图表","categories":["A","B"],"series":[{"name":"S1","values":[1,2]}]}, ...]
    }
  ],
  "citations": [{"index":1,"title":"来源","url":"https://..."}, ...]
}
"""

import json
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ── Constants ──

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_FAMILY = 'Microsoft YaHei'

# ── Helpers ──

def add_text_box(slide, text, x, y, w, h, font_size=14, bold=False, color="333333", align=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor.from_string(color)
    p.font.name = FONT_FAMILY
    p.alignment = align
    return txBox

def add_table_shape(slide, data, x, y, w):
    """添加表格。data[0] 是表头。"""
    rows = len(data)
    cols = max(len(row) for row in data)
    row_h = Inches(0.28)
    h = row_h * rows

    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), h)
    table = table_shape.table

    # Set column widths evenly
    col_w = Inches(w / cols)
    for c in range(cols):
        table.columns[c].width = col_w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if c < len(data[r]) else ""

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8) if r > 0 else Pt(9)
                p.font.name = FONT_FAMILY
                p.alignment = PP_ALIGN.CENTER

            if r == 0:
                # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string("16213E")
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string("F8F9FA")
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor.from_string("333333")

    return h

def add_chart_shape(slide, chart_spec, x, y, w, h):
    """添加图表"""
    chart_type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "line": XL_CHART_TYPE.LINE_MARKERS,
    }

    chart_type = chart_type_map.get(chart_spec.get("type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = chart_spec.get("categories", [])

    for series in chart_spec.get("series", []):
        name = series.get("name", "")
        values = series.get("values", [])
        # Scatter 特殊处理
        if chart_spec.get("type") == "scatter":
            if values and isinstance(values[0], list):
                # [[x,y], ...] → 用 y 作为值
                values = [v[1] if len(v) > 1 else 0 for v in values]
        chart_data.add_series(name, values)

    chart_shape = slide.shapes.add_chart(
        chart_type, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    )
    chart = chart_shape.chart

    # Title
    title_text = chart_spec.get("title", "")
    if title_text:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text
        for p in chart.chart_title.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.name = FONT_FAMILY

    # Legend
    chart.has_legend = len(chart_spec.get("series", [])) > 1

    return chart_shape

# ── Main ──

def generate_pptx(input_path, output_path):
    with open(input_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    title = data.get("title", "Untitled")
    sections = data.get("sections", [])
    citations = data.get("citations", [])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ═══ Title Slide ═══
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    add_text_box(slide, title, 0.5, 2.0, 12, 1.5, font_size=44, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)

    from datetime import datetime
    date_str = datetime.now().strftime("%Y/%-m/%-d")
    add_text_box(slide, date_str, 0.5, 3.8, 12, 0.5, font_size=18, color="CCCCCC", align=PP_ALIGN.CENTER)

    # ═══ Content Slides ═══
    for section in sections:
        s_title = section.get("title", "")
        tables = section.get("tables", [])
        chart_specs = section.get("chartSpecs", [])

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        # Section title
        add_text_box(slide, s_title, 0.5, 0.3, 12, 0.8, font_size=28, bold=True, color="1A1A2E")

        # Separator line
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0.5), Inches(1.1), Inches(9), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string("16213E")
        line.line.fill.background()

        # Parse HTML content to extract headings (h3) and paragraphs (p)
        import re
        html = section.get("content", "")
        html = re.sub(r'<script[\s\S]*?</script>', '', html)
        html = re.sub(r'<table[\s\S]*?</table>', '', html)
        html = re.sub(r'<sup>[\s\S]*?</sup>', '', html)

        content_parts = []
        for m in re.finditer(r'<(h3|p)\b[^>]*>(.*?)</\1>', html, re.DOTALL):
            tag = m.group(1)
            text = re.sub(r'<[^>]+>', '', m.group(2)).strip()
            if text:
                content_parts.append({'tag': tag, 'text': text})

        # Render text, tracking used height
        text_y = 1.35
        for part in content_parts[:12]:
            if part['tag'] == 'h3':
                h = max(0.28, len(part['text']) / 45 * 0.28)
                add_text_box(slide, part['text'], 0.5, text_y, 12, h, font_size=13, bold=True, color="1A1A2E")
                text_y += h + 0.05
            else:
                h = max(0.22, len(part['text']) / 65 * 0.24)
                add_text_box(slide, part['text'], 0.5, text_y, 12, h, font_size=10, color="333333")
                text_y += h + 0.05

        # Calculate remaining space for tables + charts
        y_start = text_y + 0.08
        available_h = 7.0 - y_start
        pair_count = max(len(tables), len(chart_specs))
        zone_h = min(1.3, available_h / max(pair_count, 1))

        for p_idx in range(pair_count):
            table_data = tables[p_idx] if p_idx < len(tables) else None
            chart_spec = chart_specs[p_idx] if p_idx < len(chart_specs) else None

            cy = y_start + p_idx * zone_h

            # Table on left
            if table_data:
                add_table_shape(slide, table_data, 0.5, cy, 5.5)

            # Chart on right
            if chart_spec:
                add_chart_shape(slide, chart_spec, 6.3, cy, 5.5, zone_h - 0.1)

    # ═══ Citations Slide ═══
    if citations:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text_box(slide, "参考来源", 0.5, 0.3, 12, 0.8, font_size=28, bold=True, color="1A1A2E")

        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string("16213E")
        line.line.fill.background()

        cite_text = ""
        for c in citations:
            idx = c.get("index", "")
            c_title = c.get("title", "")
            url = c.get("url", "")
            line_text = f"[{idx}] {c_title}"
            if url:
                line_text += f" {url}"
            cite_text += line_text + "\n"

        add_text_box(slide, cite_text, 0.5, 1.4, 12, 4.5, font_size=14, color="333333")

    # ── Save ──
    prs.save(output_path)
    return len(prs.slides)

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <input.json> <output.pptx>")
        sys.exit(1)

    count = generate_pptx(sys.argv[1], sys.argv[2])
    print(f"OK: {count} slides -> {os.path.getsize(sys.argv[2])} bytes")
